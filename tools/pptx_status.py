#!/usr/bin/env python3
"""Simple PPTX generator for Status slides from Markdown-like definitions.

This is a minimal helper. Requires python-pptx:
  pip install python-pptx
"""
import sys
try:
    from pptx import Presentation
    from pptx.util import Inches
except Exception:
    print("python-pptx is not installed. Install with: pip install python-pptx")
    sys.exit(1)

def create_basic_pptx(slides, out_path):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    for idx, slide in enumerate(slides):
        if idx == 0:
            slide_obj = prs.slides.add_slide(title_slide_layout)
            slide_obj.shapes.title.text = slide
        else:
            blank = prs.slide_layouts[6]
            slide_obj = prs.slides.add_slide(blank)
            txBox = slide_obj.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8.5), Inches(0.8))
            tf = txBox.text_frame
            tf.text = slide
    prs.save(out_path)
    print(f"Wrote {out_path}")

def main():
    slides = [
        "Status: SkillGap AI Pro – Enterprise Career Predictor",
        "Data & Datasets: overview of inputs and cleanup",
        "Models: RF & DT with feature engineering (~75% accuracy)",
        "APIs & UI: /api/analyze, /api/progress, UI features",
        "PDF Resume Errors: Scanned/PWD handling and user messages",
        "Next Steps: NLP-lite vs DL-based matching",
    ]
    create_basic_pptx(slides, 'presentation/Status_SkillGap_AI_Pro.pptx')

if __name__ == '__main__':
    main()
